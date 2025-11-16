from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

DEPT_COLOR = {
    '校办工作': RGBColor(255, 255, 255),      # 白色
    '教务处工作':  RGBColor(246, 208, 80),       # 主橙
    '学生处工作': RGBColor(255, 255, 255),      # 红色
    '团委工作': RGBColor(246, 208, 80),     # 红色
    '艺体中心工作': RGBColor(255, 255, 255),       # 红色
    '品招中心工作':  RGBColor(246, 208, 80),       # 红色
    '总务处、餐厅工作': RGBColor(255, 255, 255),      # 红色
    '高三年级组工作':  RGBColor(246, 208, 80),       # 红色
    '高二年级组工作': RGBColor(255, 255, 255),       # 红色
    '高一年级组工作': RGBColor(246, 208, 80),      # 红色

}
def copy_run(s, t):
    if s and t:
        t.text = s.text
        t.font.name = s.font.name
        t.font.bold = s.font.bold
        t.font.size = s.font.size
        t.font.italic = s.font.italic
        t.font.underline = s.font.underline


def copy_shape_to_slide(source_shape, target_slide, calendar):
    # 获取形状的位置和尺寸
    left = source_shape.left
    top = source_shape.top
    width = source_shape.width
    height = source_shape.height

    if source_shape.has_table:
        # 复制表格
        table = source_shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        for shape in target_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape.fill.color.rgb = RGBColor(255, 255, 255)  # 设置填充颜色为白色
                # shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 设置填充颜色为红色
                break

        new_table = target_slide.shapes.add_table(rows, cols, left, top, width, height).table

        for i in range(rows) :
            new_table.rows[i].height = table.rows[i].height

        for j in range(cols):
            new_table.columns[j].width = table.columns[j].width


        # 复制表格内容和格式
        for i in range(rows):
            for j in range(cols):
                source_cell = table.cell(i, j)
                target_cell = new_table.cell(i, j)
                target_cell.vertical_anchor = source_cell.vertical_anchor

                source_run = None
                target_run = None
                target_paragraph = target_cell.text_frame.paragraphs[0]
                source_paragraph = source_cell.text_frame.paragraphs[0]
                target_paragraph.alignment = source_paragraph.alignment

                if target_paragraph.runs:
                    target_run = target_paragraph.runs[0]
                else:
                    target_run = target_paragraph.add_run()

                if source_paragraph.runs:
                    source_run = source_paragraph.runs[0]

                copy_run(source_run, target_run)

                if i == 0 and j == 3 and target_run:
                    dept = calendar['dept']
                    target_run.text = dept
                    target_run.font.color.rgb = DEPT_COLOR[dept]

                if i == 1 and j == 0 and target_run:
                    target_run.text = calendar['week']

                if i == 1 and j == 1 and target_run:
                    target_run.text = calendar['date']

                if i == 1 and j == 2 and target_run:
                    target_run.text = calendar['weekday']

                if i == 1 and j == 3 and target_run:
                    target_run.text = calendar['content']
                    text_len = len(target_run.text) / 2

                    if text_len < 40:
                        target_run.font.size = Pt(60)
                    elif 40 <= text_len < 100:
                        target_run.font.size = Pt(40)
                    elif 100 <= text_len < 200:
                        target_run.font.size = Pt(30)
                    elif 200 <= text_len < 300:
                        target_run.font.size = Pt(20)


                # 复制文本格式
                # if source_cell.text_frame.paragraphs:
                #     source_paragraph = source_cell.text_frame.paragraphs[0]
                #     if source_paragraph.runs:
                        # source_run = source_paragraph.runs[0]

                        # if j == 3:
                        #     target_paragraph.alignment = PP_ALIGN.CENTER

                        # if j in (0, 1, 2):
                        #     target_paragraph.alignment = MSO_VERTICAL_ANCHOR.MIDDLE

                        # if target_paragraph.runs:
                        #     for target_run in target_paragraph.runs:
                        #         # target_run = target_paragraph.runs[0]
                        #         target_run.font.name = source_run.font.name
                        #         target_run.font.bold = source_run.font.bold
                        #         target_run.font.size = source_run.font.size
                        #         if i == 1 and j == 3:
                        #             target_run.font.size = Pt(40)


def fill_ppt(calendars, target_ppt_path):
    # 加载PPT模板
    # print(f"正在加载PPT模板: {'./src/2.pptx'}")
    source_ppt = Presentation('./src/2.pptx')

    # 选择要复制的幻灯片索引（例如，索引为0的幻灯片）
    slide_to_copy = source_ppt.slides[0]

    # 加载目标PPT文件或创建一个新的PPT文件
    # target_ppt_path =

    try:
        target_ppt = Presentation('./src/2.pptx')  # 如果文件已存在，则加载它
    except FileNotFoundError:
        target_ppt = Presentation()  # 如果文件不存在，创建一个新的PPT文件

    # 将选中的幻灯片复制到目标PPT中
    # target_slide_layout = slide_to_copy.slide_layout.slide_layout_id
    # new_slide = target_ppt.slides.add_slide(slide_to_copy.slide_layout)

    # 删除第一张幻灯片
    # 获取第一张幻灯片的关联ID
    first_slide_rId = target_ppt.slides._sldIdLst[0].rId

    # 删除关系
    target_ppt.part.drop_rel(first_slide_rId)

    # 从幻灯片ID列表中移除第一项
    del target_ppt.slides._sldIdLst[0]

    # shape_len = len(slide_to_copy.shapes)
    # 复制所有形状
    for calendar in calendars:
        new_slide = target_ppt.slides.add_slide(slide_to_copy.slide_layout)
        for shape in slide_to_copy.shapes:
            copy_shape_to_slide(shape, new_slide, calendar)

    print('输出文件地址：${}', target_ppt_path)
    # 保存目标PPT文件
    target_ppt.save(target_ppt_path)
